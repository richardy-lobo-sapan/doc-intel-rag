"""
PPTX extraction via python-pptx.

Granularity: one RawSegment per slide. We concatenate all text frames on
a slide (title + body + any text boxes) in shape order, and separately
pull speaker notes if present since they often contain procedural detail
not visible on the slide itself.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.extractors.base import BaseExtractor, RawSegment


class PPTXExtractor(BaseExtractor):
    supported_extensions = (".pptx",)

    def extract(self, file_path: Path) -> list[RawSegment]:
        prs = Presentation(str(file_path))
        segments: list[RawSegment] = []
        total_slides = len(prs.slides)

        for slide_index, slide in enumerate(prs.slides, start=1):
            text_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_text = "\n".join(
                        p.text for p in shape.text_frame.paragraphs if p.text.strip()
                    )
                    if shape_text.strip():
                        text_parts.append(shape_text)

            notes_text = ""
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()

            slide_text = "\n".join(text_parts).strip()
            if notes_text:
                slide_text = f"{slide_text}\n[Speaker notes: {notes_text}]"

            if slide_text.strip():
                segments.append(
                    RawSegment(
                        text=slide_text,
                        source_file=file_path.name,
                        location_type="slide",
                        location_value=str(slide_index),
                        extra={"total_slides": total_slides},
                    )
                )

        return segments
