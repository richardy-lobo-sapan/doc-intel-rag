"""
Generates a small fictional internal-document corpus (SOPs, audit
procedures, compliance reports) across all supported formats, used to
exercise and demo the extraction + RAG pipeline. This is sample/test data
only -- not real organizational documents.
"""

import csv
from pathlib import Path

import docx
import openpyxl
from docx.shared import Pt
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.platypus import (
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

OUT_DIR = Path(__file__).resolve().parent.parent / "data" / "corpus"
OUT_DIR.mkdir(parents=True, exist_ok=True)


# ---------------------------------------------------------------------------
# PDF documents
# ---------------------------------------------------------------------------

def make_pdf_audit_procedure():
    path = OUT_DIR / "sop_quality_audit_production_line.pdf"
    doc = SimpleDocTemplate(str(path), pagesize=A4, topMargin=2 * cm, bottomMargin=2 * cm)
    styles = getSampleStyleSheet()
    h1 = ParagraphStyle("H1", parent=styles["Heading1"], fontSize=16, spaceAfter=12)
    h2 = ParagraphStyle("H2", parent=styles["Heading2"], fontSize=13, spaceAfter=8)
    body = styles["BodyText"]

    elems = []
    elems.append(Paragraph("Standard Operating Procedure: Quality Audit for Production Line B", h1))
    elems.append(Paragraph("Document No: SOP-QA-014 | Revision: 3 | Effective Date: 2026-01-15", body))
    elems.append(Spacer(1, 12))

    elems.append(Paragraph("1. Purpose", h2))
    elems.append(Paragraph(
        "This procedure defines the steps required to conduct a quality audit on "
        "Production Line B, covering raw material inspection, in-process checks, "
        "and final product verification before release to the warehouse.", body))
    elems.append(Spacer(1, 8))

    elems.append(Paragraph("2. Scope", h2))
    elems.append(Paragraph(
        "Applies to all Quality Assurance auditors and Line B production supervisors. "
        "Excludes Production Lines A and C, which follow SOP-QA-012 and SOP-QA-013 respectively.", body))
    elems.append(Spacer(1, 8))

    elems.append(Paragraph("3. Audit Frequency", h2))
    elems.append(Paragraph(
        "Routine audits are conducted weekly, every Monday at 08:00. Unscheduled "
        "audits may be triggered by a customer complaint, a non-conformance report (NCR), "
        "or a request from the Plant Manager.", body))

    elems.append(Spacer(1, 16))
    elems.append(Paragraph("4. Audit Procedure Steps", h2))
    elems.append(Paragraph(
        "Step 1: Verify that the production line log matches the master schedule for "
        "the current shift. Any discrepancy must be recorded as a deviation.", body))
    elems.append(Spacer(1, 6))
    elems.append(Paragraph(
        "Step 2: Inspect a random sample of 15 units from the current batch using the "
        "dimensional checklist QA-FORM-09. Reject the batch if more than 2 units fail.", body))
    elems.append(Spacer(1, 6))
    elems.append(Paragraph(
        "Step 3: Check that all operators on shift have a valid and current safety "
        "certification on file. Flag any expired certification to HR within 24 hours.", body))
    elems.append(Spacer(1, 6))
    elems.append(Paragraph(
        "Step 4: Review the calibration sticker on all measurement tools used on the "
        "line. Tools with calibration overdue by more than 7 days must be pulled from "
        "service immediately and sent to the metrology lab.", body))
    elems.append(Spacer(1, 6))
    elems.append(Paragraph(
        "Step 5: Document all findings in the Audit Report Form QA-FORM-11 and submit "
        "to the QA Manager within 4 hours of audit completion.", body))

    elems.append(Spacer(1, 16))
    elems.append(Paragraph("5. Responsibility Matrix", h2))
    table_data = [
        ["Role", "Responsibility"],
        ["QA Auditor", "Conduct audit, complete QA-FORM-11"],
        ["Line Supervisor", "Provide access, production logs"],
        ["QA Manager", "Review findings, approve/reject batch"],
        ["Plant Manager", "Escalation point for repeated failures"],
    ]
    table = Table(table_data, colWidths=[6 * cm, 10 * cm])
    table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2c3e50")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 6),
    ]))
    elems.append(table)

    elems.append(Spacer(1, 16))
    elems.append(Paragraph("6. Non-Conformance Handling", h2))
    elems.append(Paragraph(
        "If a batch is rejected during audit, the Line Supervisor must quarantine "
        "the affected units within 1 hour and initiate NCR-FORM-02. Production may "
        "resume only after QA Manager sign-off on corrective action.", body))

    doc.build(elems)
    return path.name


def make_pdf_data_retention_policy():
    path = OUT_DIR / "regulation_data_retention_policy.pdf"
    doc = SimpleDocTemplate(str(path), pagesize=A4, topMargin=2 * cm, bottomMargin=2 * cm)
    styles = getSampleStyleSheet()
    h1 = ParagraphStyle("H1", parent=styles["Heading1"], fontSize=16, spaceAfter=12)
    h2 = ParagraphStyle("H2", parent=styles["Heading2"], fontSize=13, spaceAfter=8)
    body = styles["BodyText"]

    elems = []
    elems.append(Paragraph("Internal Regulation: Data Retention and Disposal Policy", h1))
    elems.append(Paragraph("Document No: REG-IT-007 | Revision: 5 | Effective Date: 2025-11-01", body))
    elems.append(Spacer(1, 12))

    elems.append(Paragraph("1. Policy Statement", h2))
    elems.append(Paragraph(
        "All business records, whether physical or digital, must be retained only "
        "for the period required by applicable law or operational need, and securely "
        "disposed of thereafter to reduce legal and security risk.", body))

    elems.append(Spacer(1, 8))
    elems.append(Paragraph("2. Retention Periods by Record Type", h2))
    elems.append(Paragraph(
        "Financial records (invoices, tax filings): 10 years. "
        "Employee records: 5 years after termination. "
        "Customer support tickets: 3 years. "
        "Internal audit reports: 7 years. "
        "Marketing materials: 2 years or until superseded.", body))

    elems.append(Spacer(1, 8))
    elems.append(Paragraph("3. Disposal Procedure", h2))
    elems.append(Paragraph(
        "Physical documents past their retention period must be shredded using a "
        "cross-cut shredder rated P-4 or higher. Digital records must be deleted using "
        "a secure-erase utility and the deletion logged in the IT Asset Management "
        "system with a timestamp and the responsible employee's ID.", body))

    elems.append(Spacer(1, 8))
    elems.append(Paragraph("4. Exceptions", h2))
    elems.append(Paragraph(
        "Records subject to an active legal hold must NOT be disposed of regardless "
        "of retention period expiry. Legal holds are issued by the Legal department "
        "and override this policy until formally lifted in writing.", body))

    elems.append(Spacer(1, 8))
    elems.append(Paragraph("5. Roles and Accountability", h2))
    elems.append(Paragraph(
        "Department heads are accountable for ensuring records under their control "
        "follow this policy. The IT Security team audits compliance quarterly and "
        "reports findings to the Compliance Committee.", body))

    elems.append(Spacer(1, 8))
    elems.append(Paragraph("6. Violations", h2))
    elems.append(Paragraph(
        "Failure to comply with this policy, including premature disposal of records "
        "under legal hold, is treated as a serious compliance violation and may result "
        "in disciplinary action up to and including termination.", body))

    doc.build(elems)
    return path.name


# ---------------------------------------------------------------------------
# DOCX documents
# ---------------------------------------------------------------------------

def make_docx_incident_response():
    path = OUT_DIR / "sop_incident_response_it_security.docx"
    document = docx.Document()

    document.add_heading("SOP: IT Security Incident Response", level=0)
    p = document.add_paragraph("Document No: SOP-SEC-021 | Revision: 2 | Effective Date: 2026-02-01")
    p.runs[0].font.size = Pt(10)

    document.add_heading("1. Purpose", level=1)
    document.add_paragraph(
        "This procedure defines the steps the IT Security team must follow when "
        "responding to a suspected or confirmed security incident, from initial "
        "detection through to post-incident review."
    )

    document.add_heading("2. Incident Severity Levels", level=1)
    document.add_paragraph(
        "Severity 1 (Critical): Active data breach or ransomware affecting production "
        "systems. Response required within 15 minutes.\n"
        "Severity 2 (High): Suspicious activity on a single system, no confirmed breach. "
        "Response required within 1 hour.\n"
        "Severity 3 (Medium): Phishing report or policy violation with no system impact. "
        "Response required within 4 business hours."
    )

    document.add_heading("3. Detection and Triage", level=1)
    document.add_paragraph(
        "Upon receiving an alert from the SIEM system or a report from an employee, "
        "the on-call security analyst must triage within the response window defined "
        "by severity level. Triage includes confirming the incident is genuine, "
        "assigning a severity level, and opening a ticket in the incident tracker."
    )

    document.add_heading("4. Containment", level=1)
    document.add_paragraph(
        "For Severity 1 incidents, the affected system must be isolated from the "
        "network immediately, even before full investigation, to prevent lateral "
        "movement. Isolation requires sign-off from the Security Lead or, if "
        "unavailable, the IT Director."
    )

    document.add_heading("5. Communication Protocol", level=1)
    document.add_paragraph(
        "Severity 1 incidents require notifying the CISO within 30 minutes and "
        "Legal within 2 hours, due to potential regulatory disclosure obligations. "
        "All external communication about an incident must be approved by Legal "
        "and Corporate Communications before release."
    )

    document.add_heading("6. Recovery and Closure", level=1)
    document.add_paragraph(
        "Systems may be returned to production only after the Security Lead confirms "
        "the root cause has been remediated and a vulnerability scan shows no "
        "remaining indicators of compromise. The incident ticket is closed only after "
        "a post-incident review document is filed."
    )

    document.add_heading("7. Roles", level=1)
    table = document.add_table(rows=1, cols=2)
    table.style = "Light Grid Accent 1"
    hdr = table.rows[0].cells
    hdr[0].text = "Role"
    hdr[1].text = "Responsibility"
    rows_data = [
        ("On-call Security Analyst", "First responder, triage, initial containment"),
        ("Security Lead", "Approves containment actions, coordinates response"),
        ("CISO", "Notified for Sev-1, makes disclosure decisions with Legal"),
        ("IT Director", "Backup approver when Security Lead unavailable"),
    ]
    for role, resp in rows_data:
        row = table.add_row().cells
        row[0].text = role
        row[1].text = resp

    document.save(str(path))
    return path.name


def make_docx_procurement_policy():
    path = OUT_DIR / "regulation_procurement_approval_policy.docx"
    document = docx.Document()

    document.add_heading("Internal Regulation: Procurement Approval Policy", level=0)
    p = document.add_paragraph("Document No: REG-FIN-009 | Revision: 4 | Effective Date: 2025-09-01")
    p.runs[0].font.size = Pt(10)

    document.add_heading("1. Policy Statement", level=1)
    document.add_paragraph(
        "All procurement of goods and services must follow a tiered approval process "
        "based on transaction value, to ensure appropriate financial oversight and "
        "prevent unauthorized spending."
    )

    document.add_heading("2. Approval Thresholds", level=1)
    document.add_paragraph(
        "Purchases under IDR 5,000,000 may be approved by the requesting Department "
        "Head directly.\n"
        "Purchases between IDR 5,000,000 and IDR 50,000,000 require approval from the "
        "Finance Manager in addition to the Department Head.\n"
        "Purchases above IDR 50,000,000 require approval from the CFO and must be "
        "accompanied by at least two competing vendor quotations."
    )

    document.add_heading("3. Vendor Selection", level=1)
    document.add_paragraph(
        "Preferred vendors on the approved vendor list may be used without additional "
        "quotations for purchases below IDR 50,000,000. New vendors must complete a "
        "due diligence questionnaire and be approved by Procurement before any purchase "
        "order is issued."
    )

    document.add_heading("4. Emergency Procurement", level=1)
    document.add_paragraph(
        "In situations posing immediate risk to safety or operations, a Department "
        "Head may authorize emergency procurement up to IDR 20,000,000 without prior "
        "Finance Manager approval, but must file a retroactive justification within "
        "48 hours."
    )

    document.add_heading("5. Conflict of Interest", level=1)
    document.add_paragraph(
        "Any employee with a personal or financial relationship with a vendor under "
        "consideration must disclose the relationship to Procurement before the "
        "purchase decision is made. Failure to disclose is grounds for disciplinary "
        "action regardless of whether the vendor was ultimately selected."
    )

    document.add_heading("6. Audit Trail Requirements", level=1)
    document.add_paragraph(
        "Every purchase order, regardless of value, must be logged in the procurement "
        "system with the requester, approver(s), vendor, amount, and business "
        "justification. Records must be retained per REG-IT-007 (Data Retention Policy)."
    )

    document.save(str(path))
    return path.name


# ---------------------------------------------------------------------------
# PPTX documents
# ---------------------------------------------------------------------------

def _add_bullet_slide(prs, title, bullets, notes=None):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = bullets[0]
    for b in bullets[1:]:
        p = tf.add_paragraph()
        p.text = b
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def make_pptx_onboarding_training():
    path = OUT_DIR / "training_new_employee_safety_onboarding.pptx"
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "New Employee Safety Onboarding"
    title_slide.placeholders[1].text = "Training Deck v3 - Effective 2026-01-01"

    _add_bullet_slide(
        prs, "Why Safety Onboarding Matters",
        [
            "All new employees must complete this training before entering production areas",
            "Required by REG-HSE-003 within first 3 working days of employment",
            "Completion is logged in the HR Learning Management System",
        ],
        notes="Trainer note: confirm attendance sheet is signed before starting slide 3.",
    )

    _add_bullet_slide(
        prs, "Personal Protective Equipment (PPE)",
        [
            "Safety glasses and steel-toe boots are mandatory in all production zones",
            "Hearing protection required in zones marked with the orange noise symbol",
            "PPE is issued at the warehouse counter on Day 1 -- report missing PPE to your supervisor immediately",
        ],
    )

    _add_bullet_slide(
        prs, "Emergency Evacuation Procedure",
        [
            "On hearing the continuous siren, stop work immediately and proceed to the nearest marked exit",
            "Assemble at Muster Point A (north parking lot) or Muster Point B (south gate)",
            "Do not re-enter the building until the Fire Warden gives the all-clear signal",
        ],
        notes="Trainer note: walk new hires to their nearest muster point physically on Day 1.",
    )

    _add_bullet_slide(
        prs, "Reporting Near-Misses and Incidents",
        [
            "Any near-miss, however minor, must be reported using the QR code posted at each workstation",
            "Reports go directly to the HSE team and are reviewed within 24 hours",
            "Employees will never face disciplinary action for reporting in good faith",
        ],
    )

    _add_bullet_slide(
        prs, "Sign-Off and Certification",
        [
            "Complete the online quiz (minimum 80% to pass) within 24 hours of this session",
            "Certification is valid for 12 months and must be renewed annually",
            "Questions? Contact the HSE team at hse-support internal extension 4420",
        ],
    )

    prs.save(str(path))
    return path.name


def make_pptx_quarterly_compliance_review():
    path = OUT_DIR / "report_quarterly_compliance_review_q1.pptx"
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Quarterly Compliance Review - Q1 2026"
    title_slide.placeholders[1].text = "Compliance Committee Presentation"

    _add_bullet_slide(
        prs, "Audit Summary",
        [
            "14 scheduled quality audits conducted across Production Lines A, B, and C",
            "2 audits resulted in batch rejection, both on Line B due to dimensional non-conformance",
            "Average audit completion time: 2.3 hours, within the 4-hour SOP target",
        ],
    )

    _add_bullet_slide(
        prs, "Security Incidents This Quarter",
        [
            "3 Severity 3 incidents (phishing reports), all contained within SLA",
            "0 Severity 1 or 2 incidents reported",
            "Average phishing report response time: 1.8 hours, against a 4-hour target",
        ],
    )

    _add_bullet_slide(
        prs, "Procurement Policy Compliance",
        [
            "All purchase orders above IDR 50,000,000 included two vendor quotations as required",
            "1 emergency procurement case filed, justification submitted within the 48-hour window",
            "No conflict-of-interest disclosures filed this quarter",
        ],
    )

    _add_bullet_slide(
        prs, "Data Retention Compliance",
        [
            "Quarterly digital records audit found 99.2% compliance with disposal timelines",
            "2 records flagged for premature disposal review, both cleared as not under legal hold",
            "Next IT Security quarterly audit scheduled for the second week of next quarter",
        ],
    )

    _add_bullet_slide(
        prs, "Action Items for Next Quarter",
        [
            "Review Line B calibration schedule following the 2 dimensional non-conformance rejections",
            "Refresh new-employee safety onboarding deck with updated muster point map",
            "Pilot automated purchase-order quotation tracking in the procurement system",
        ],
    )

    prs.save(str(path))
    return path.name


# ---------------------------------------------------------------------------
# XLSX documents
# ---------------------------------------------------------------------------

def make_xlsx_audit_log():
    path = OUT_DIR / "log_quality_audit_records_2026.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Audit Log"
    ws.append(["Audit Date", "Production Line", "Auditor", "Sample Size", "Units Rejected", "Result", "Notes"])
    rows = [
        ["2026-01-05", "Line A", "Siti Rahma", 15, 0, "Pass", "No issues found"],
        ["2026-01-05", "Line B", "Budi Santoso", 15, 1, "Pass", "Minor surface scratch, within tolerance"],
        ["2026-01-12", "Line B", "Budi Santoso", 15, 3, "Fail", "Dimensional non-conformance, batch quarantined"],
        ["2026-01-12", "Line C", "Siti Rahma", 15, 0, "Pass", "No issues found"],
        ["2026-01-19", "Line A", "Andi Wijaya", 15, 0, "Pass", "No issues found"],
        ["2026-01-19", "Line B", "Budi Santoso", 15, 4, "Fail", "Dimensional non-conformance again, escalated to QA Manager"],
        ["2026-01-26", "Line C", "Siti Rahma", 15, 0, "Pass", "No issues found"],
        ["2026-02-02", "Line A", "Andi Wijaya", 15, 1, "Pass", "Within tolerance"],
        ["2026-02-02", "Line B", "Budi Santoso", 15, 0, "Pass", "Calibration recheck resolved prior issue"],
        ["2026-02-09", "Line C", "Siti Rahma", 15, 0, "Pass", "No issues found"],
    ]
    for row in rows:
        ws.append(row)

    ws2 = wb.create_sheet("Calibration Status")
    ws2.append(["Tool ID", "Tool Type", "Line", "Last Calibrated", "Next Due", "Status"])
    cal_rows = [
        ["CAL-101", "Digital Caliper", "Line A", "2025-12-01", "2026-06-01", "Valid"],
        ["CAL-102", "Digital Caliper", "Line B", "2025-10-15", "2026-04-15", "Valid"],
        ["CAL-103", "Torque Wrench", "Line B", "2025-09-01", "2026-03-01", "Overdue"],
        ["CAL-104", "Digital Caliper", "Line C", "2025-12-20", "2026-06-20", "Valid"],
        ["CAL-105", "Height Gauge", "Line A", "2025-11-10", "2026-05-10", "Valid"],
    ]
    for row in cal_rows:
        ws2.append(row)

    wb.save(str(path))
    return path.name


def make_xlsx_security_incident_tracker():
    path = OUT_DIR / "log_security_incident_tracker_2026.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Incidents"
    ws.append(["Incident ID", "Date Reported", "Severity", "Type", "Response Time (hrs)", "Status", "Resolved By"])
    rows = [
        ["INC-2026-001", "2026-01-08", "Sev 3", "Phishing report", 1.5, "Closed", "Analyst: Rina"],
        ["INC-2026-002", "2026-01-22", "Sev 3", "Phishing report", 2.1, "Closed", "Analyst: Rina"],
        ["INC-2026-003", "2026-02-03", "Sev 3", "Suspicious login attempt", 1.9, "Closed", "Analyst: Fajar"],
        ["INC-2026-004", "2026-02-18", "Sev 2", "Malware alert on workstation", 0.8, "Closed", "Security Lead: Dewi"],
    ]
    for row in rows:
        ws.append(row)

    ws2 = wb.create_sheet("Response SLA Targets")
    ws2.append(["Severity", "Target Response Time", "Notification Required"])
    sla_rows = [
        ["Sev 1", "15 minutes", "CISO within 30 min, Legal within 2 hrs"],
        ["Sev 2", "1 hour", "Security Lead"],
        ["Sev 3", "4 business hours", "None required"],
    ]
    for row in sla_rows:
        ws2.append(row)

    wb.save(str(path))
    return path.name


# ---------------------------------------------------------------------------
# CSV documents
# ---------------------------------------------------------------------------

def make_csv_vendor_list():
    path = OUT_DIR / "approved_vendor_list.csv"
    with open(path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["Vendor Name", "Category", "Approval Date", "Max Purchase Without Quote (IDR)", "Status"])
        writer.writerows([
            ["PT Sumber Makmur Logistik", "Logistics", "2024-03-01", "50000000", "Active"],
            ["PT Cipta Mandiri Supplies", "Office Supplies", "2023-11-15", "50000000", "Active"],
            ["PT Baja Presisi Indonesia", "Raw Materials", "2025-01-10", "50000000", "Active"],
            ["PT Teknologi Aman Sentosa", "IT Equipment", "2025-06-01", "50000000", "Active"],
            ["PT Karya Bersama Konstruksi", "Facilities", "2022-08-20", "50000000", "Under Review"],
            ["CV Sinar Jaya Print", "Printing Services", "2024-09-05", "50000000", "Active"],
        ])
    return path.name


def make_csv_safety_certification_records():
    path = OUT_DIR / "employee_safety_certification_records.csv"
    with open(path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["Employee ID", "Name", "Line Assignment", "Certification Date", "Expiry Date", "Status"])
        writer.writerows([
            ["EMP-0231", "Andi Wijaya", "Line A", "2025-06-01", "2026-06-01", "Valid"],
            ["EMP-0245", "Budi Santoso", "Line B", "2025-03-15", "2026-03-15", "Valid"],
            ["EMP-0198", "Siti Rahma", "Line C", "2025-08-22", "2026-08-22", "Valid"],
            ["EMP-0301", "Dewi Lestari", "Line B", "2024-12-01", "2025-12-01", "Expired"],
            ["EMP-0312", "Fajar Hidayat", "Line A", "2025-09-10", "2026-09-10", "Valid"],
            ["EMP-0299", "Rina Wulandari", "IT Security", "2025-05-01", "2026-05-01", "Valid"],
        ])
    return path.name


def main():
    generated = []
    generated.append(make_pdf_audit_procedure())
    generated.append(make_pdf_data_retention_policy())
    generated.append(make_docx_incident_response())
    generated.append(make_docx_procurement_policy())
    generated.append(make_pptx_onboarding_training())
    generated.append(make_pptx_quarterly_compliance_review())
    generated.append(make_xlsx_audit_log())
    generated.append(make_xlsx_security_incident_tracker())
    generated.append(make_csv_vendor_list())
    generated.append(make_csv_safety_certification_records())

    print(f"Generated {len(generated)} files in {OUT_DIR}:")
    for name in generated:
        print(f"  - {name}")


if __name__ == "__main__":
    main()
